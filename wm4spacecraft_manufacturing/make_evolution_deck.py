#!/usr/bin/env python
"""Build ConstructionBots_evolution.pptx — what changed from SISL ConstructionBots to this repo.

Design language reused from ConstructionBots_0706_revised.pptx (16:9, navy/orange/teal).
Run:  python make_evolution_deck.py
"""
# =====================================================================
# [이 파일이 하는 일] (초보자용 한국어 설명)
# ---------------------------------------------------------------------
# 목적: PowerPoint 발표자료(.pptx)를 "코드로" 자동 생성하는 스크립트.
#       손으로 슬라이드를 그리지 않고, 파이썬으로 도형/텍스트를 찍어서
#       ConstructionBots_evolution.pptx 라는 16:9 발표덱을 만든다.
#
# 프로젝트 안에서의 역할:
#   - ConstructionBots.jl = 여러 로봇이 협력해 물체를 조립하는 TAMP 시뮬레이터.
#     (TAMP = Task And Motion Planning, 작업배정 + 이동경로를 함께 푸는 것)
#   - 이 스토리(SISL 원본 → 현재 repo)를 발표로 설명하는 슬라이드를 만든다.
#     즉 "코드 결과물"이 아니라 "발표 자료"를 뽑아내는 유틸리티다.
#   - makespan(전체 완료시간), MILP(정수최적화 solver), RVO(충돌회피),
#     DSL(수리 명령어 문법), ForbidZone(진입금지 구역), surrogate(대리 모델),
#     respec(재명세) 같은 도메인 용어는 슬라이드 본문에 그대로 등장한다.
#
# 실행 방법:
#   python make_evolution_deck.py
#   → 같은 폴더에 ConstructionBots_evolution.pptx 파일이 생성된다.
#   (python-pptx 패키지 필요:  pip install python-pptx)
#
# ---------------------------------------------------------------------
# [문법 참고] 이 파일에서 낯설 수 있는 파이썬 문법 모음
# ---------------------------------------------------------------------
# 1) python-pptx API 객체 계층:
#      Presentation  → 파일 하나(발표덱 전체)
#      prs.slides    → 슬라이드 목록,  add_slide()로 새 장 추가
#      slide.shapes  → 그 슬라이드 위의 도형들,  add_shape()/add_textbox()로 추가
#      shape.text_frame → 도형 안 글상자,  .paragraphs / .add_paragraph()
#      paragraph.add_run() → 같은 문단 안에서 글꼴/색이 다른 텍스트 조각(run)
#    한 마디로: 발표덱 > 슬라이드 > 도형 > 글상자 > 문단 > run 순서로 중첩된다.
#
# 2) EMU 단위:  PowerPoint 내부 좌표 단위(English Metric Unit).
#      Inches(1) = 914400 EMU,  Emu(...)로 직접 값을 줄 수도 있다.
#      Pt(...)는 글자/여백의 포인트 크기.  Inches/Pt/Emu는 EMU로 변환해 주는 헬퍼.
#
# 3) 리스트/딕셔너리 컴프리헨션 & enumerate:
#      for i, (a, b, c) in enumerate(items): ...  → 항목마다 인덱스 i와 함께
#      튜플을 한 번에 풀어(unpack)서 슬라이드 위 위치를 계산하는 패턴이 반복된다.
#
# 4) 헬퍼 클로저/함수:  slide(), tb(), line(), card(), arrow() 등은 반복되는
#      "도형 만들기" 코드를 짧게 줄인 헬퍼다. 슬라이드마다 이들을 조합해 그린다.
#
# 5) f-string 은 여기선 거의 안 쓰이고, 대신 문자열을 그대로 슬라이드에 넣는다.
#      단, str(_page[0]) 처럼 숫자를 문자열로 바꾸는 형변환은 자주 나온다.
#
# 6) _page = [1] 같은 "1칸짜리 리스트"는 함수 안에서 값을 바꿔도 밖에 반영되게
#      하려는 트릭(리스트는 가변이라 _page[0] += 1이 전역 상태처럼 동작).
# =====================================================================
from pptx import Presentation                       # 발표덱(파일) 전체를 나타내는 최상위 객체
from pptx.util import Inches, Pt, Emu                # 길이 단위 헬퍼: 인치/포인트/EMU → 내부 EMU로 변환
from pptx.dml.color import RGBColor                  # 색상: RGBColor(빨강,초록,파랑) 0~255 값
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR      # 텍스트 정렬(좌/중/우) & 세로 앵커(위/가운데)
from pptx.enum.shapes import MSO_SHAPE               # 도형 종류 상수(사각형/둥근사각형/화살표/원 등)

# ---------------------------------------------------------------- palette
# 아래는 발표덱 전체에서 재사용하는 색 팔레트(네이비/오렌지/틸 톤).
# RGBColor(0x1B, 0x2A, 0x4A) = 16진수로 (R=0x1B, G=0x2A, B=0x4A) 색을 지정.
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
INK    = RGBColor(0x22, 0x27, 0x33)
ORANGE = RGBColor(0xE2, 0x6D, 0x1F)
TEAL   = RGBColor(0x12, 0x7A, 0x6E)
BLUE   = RGBColor(0x2E, 0x5C, 0x9A)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xC9, 0xD2, 0xE0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xF7, 0xF8, 0xFA)
RED    = RGBColor(0xB3, 0x32, 0x2C)
CARD_D = RGBColor(0x24, 0x38, 0x5E)   # card on dark   / 어두운 배경 위 카드 색
CARD_L = RGBColor(0xEC, 0xEF, 0xF4)   # card on light  / 밝은 배경 위 카드 색

prs = Presentation()                 # 빈 발표덱 객체 생성(여기에 슬라이드를 계속 추가한다)
prs.slide_width  = Emu(12192000)     # 슬라이드 가로 폭 = 13.333인치(16:9 와이드)를 EMU로 지정
prs.slide_height = Emu(6858000)      # 슬라이드 세로 높이 = 7.5인치를 EMU로 지정
BLANK = prs.slide_layouts[6]         # 레이아웃 6번 = 완전 빈(blank) 레이아웃, 모든 슬라이드가 이걸 쓴다
W = 13.333                           # 인치 단위 슬라이드 폭(도형 배치 계산에 반복 사용)

_page = [1]                          # 페이지 번호 카운터. 1칸 리스트라 함수 안에서 _page[0]+=1로 증가시킴


# 새 빈 슬라이드 1장을 만들고 배경 사각형을 깔아주는 헬퍼. dark=True면 네이비 배경, 아니면 밝은 배경.
def slide(dark=False):
    s = prs.slides.add_slide(BLANK)                                            # 빈 레이아웃으로 슬라이드 1장 추가
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(7.5)) # 슬라이드 전체를 덮는 배경 사각형
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY if dark else PAPER          # 채우기 색: 다크=네이비 / 라이트=종이색
    bg.line.fill.background(); bg.shadow.inherit = False                       # 테두리·그림자 제거(깔끔한 단색 배경)
    return s


# 슬라이드 왼쪽 끝에 세로 오렌지 띠(액센트 바)를 그리는 헬퍼.
def bar(s):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), Inches(7.5))  # 폭 0.18인치 세로 전체 길이 막대
    r.fill.solid(); r.fill.fore_color.rgb = ORANGE
    r.line.fill.background(); r.shadow.inherit = False


# 글상자(textbox)를 만들고 그 안의 text_frame(글틀)을 돌려주는 헬퍼. x,y=좌상단 위치, w,h=크기(인치), anchor=세로정렬.
def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    t = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor          # 자동 줄바꿈 켜고 세로 앵커 지정
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0       # 글상자 안쪽 여백 0(꽉 채우기)
    return tf


# 글틀(tf)에 한 줄(문단)을 추가하는 핵심 헬퍼. first=True면 기존 첫 문단을 쓰고, 아니면 새 문단을 덧붙인다.
# size=글자크기(pt), color=글자색, bold/italic=굵게/기울임, align=정렬, space_after=문단 뒤 간격(pt).
def line(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         space_after=0, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()   # 첫 줄이면 이미 있는 빈 문단 재사용, 아니면 새 문단 생성
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text                          # run = 같은 서식을 공유하는 글자 조각. 여기에 실제 텍스트를 넣음
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = 'Calibri'
    return p


# 슬라이드 상단 머리말을 찍는 헬퍼: 작은 오렌지 kicker(윗줄) + 큰 title(제목). dark로 제목 글자색을 흰/네이비 전환.
def head(s, kicker, title, dark=False):
    bar(s)                                                                                # 왼쪽 오렌지 띠 먼저
    line(tb(s, 0.55, 0.30, 9.0, 0.40), kicker, 13, ORANGE, True, first=True)              # 작은 라벨(섹션 표시)
    line(tb(s, 0.55, 0.72, 12.2, 1.0), title, 33, WHITE if dark else NAVY, True, first=True)  # 큰 제목


# 슬라이드 하단 꼬리말: 선택적 문구(text) + 오른쪽 아래 페이지 번호. 부를 때마다 _page 카운터를 1 증가시킨다.
def foot(s, text=None, dark=False):
    if text:
        line(tb(s, 0.70, 6.55, 11.4, 0.6), text, 18, LIGHT if dark else GRAY, True, first=True)  # 하단 한 줄 요약
    line(tb(s, 12.20, 7.05, 1.0, 0.3), str(_page[0]), 11, GRAY, first=True)  # 페이지 번호(숫자를 문자열로 변환해 표시)
    _page[0] += 1                                                            # 다음 슬라이드를 위해 번호 증가


# 둥근 사각형 "카드"를 그리는 헬퍼. accent를 주면 카드 위쪽에 얇은 색 띠를, fill로 카드 배경색을 덮어쓸 수 있다.
def card(s, x, y, w, h, dark=False, accent=None, fill=None):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill or (CARD_D if dark else CARD_L)  # fill이 있으면 그 색, 없으면 다크/라이트 기본
    c.line.fill.background(); c.shadow.inherit = False
    try:
        c.adjustments[0] = 0.05   # 둥근 모서리 곡률(0=각짐, 클수록 둥긂). 도형에 따라 실패할 수 있어 try로 감쌈
    except Exception:
        pass
    if accent:
        a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.14))  # 카드 상단 액센트 띠
        a.fill.solid(); a.fill.fore_color.rgb = accent
        a.line.fill.background(); a.shadow.inherit = False
    return c


# 오른쪽 화살표(꺾쇠) 도형을 그리는 헬퍼. 단계 사이 흐름(A→B)을 표시할 때 사용.
def arrow(s, x, y, w=0.42, h=0.20, color=None):
    """A right-pointing chevron. Keep w >= 1.5*h or PowerPoint renders it as a blob."""
    # 주의: 폭 w를 높이 h의 1.5배 이상으로 안 잡으면 PowerPoint가 화살표를 뭉개진 덩어리로 그린다.
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color or RGBColor(0xB6, 0xBE, 0xCB)  # color 없으면 기본 회색
    a.line.fill.background(); a.shadow.inherit = False
    try:
        a.adjustments[0] = 0.45   # shaft thickness  / 화살대(몸통) 두께 비율
        a.adjustments[1] = 0.55   # head width       / 화살촉 폭 비율
    except Exception:
        pass


# 데모 영상이 들어갈 자리를 표시하는 점선 테두리 플레이스홀더 박스. label=제목, sub=부제(스크립트명·길이).
def video_ph(s, x, y, w, h, label, sub):
    """Placeholder box where the recorded demo video gets embedded."""
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x0E, 0x18, 0x2E)  # 아주 어두운 남색(영상 자리 느낌)
    c.line.color.rgb = ORANGE; c.line.width = Pt(1.5); c.line.dash_style = 4  # 오렌지 점선 테두리(dash_style=4=점선)
    c.shadow.inherit = False
    tf = tb(s, x + 0.25, y + h / 2 - 0.62, w - 0.5, 1.3, MSO_ANCHOR.MIDDLE)   # 박스 세로 중앙에 안내 문구를 얹음
    line(tf, '▶  ' + label, 19, WHITE, True, first=True, align=PP_ALIGN.CENTER)   # ▶ 재생 아이콘 + 제목
    line(tf, sub, 13, LIGHT, align=PP_ALIGN.CENTER, space_after=0)                # 실행 스크립트/길이
    line(tf, '[ embed demo video here ]', 11, ORANGE, align=PP_ALIGN.CENTER)      # 나중에 실제 mp4를 끼워 넣으라는 표시


# 슬라이드의 발표자 노트(notes)에 텍스트를 넣는 헬퍼. strip()으로 앞뒤 공백/개행을 정리한다.
def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


# ================================================================ 1 · title
# [슬라이드 1 · 타이틀] 발표 제목 표지. "닫힌세계 TAMP solver → 자기진화 surrogate world model" 한 문장 주제 제시.
s = slide(dark=True)   # 표지는 어두운(네이비) 배경
# 제목 위 짧은 오렌지 밑줄 장식(강조선)
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.90), Inches(2.05), Inches(2.40), Inches(0.09))
r.fill.solid(); r.fill.fore_color.rgb = ORANGE; r.line.fill.background(); r.shadow.inherit = False
line(tb(s, 0.90, 1.40, 11.5, 0.5), 'PROJECT EVOLUTION', 16, ORANGE, True, first=True)
tf = tb(s, 0.90, 2.30, 11.6, 2.4)
line(tf, 'From a Closed-World TAMP Solver', 38, WHITE, True, first=True)
line(tf, 'to a Self-Evolving Surrogate World Model', 38, ORANGE, True)
line(tf, 'for OOD-Adaptive Spacecraft Assembly', 38, WHITE, True)
line(tb(s, 0.90, 5.15, 11.5, 0.5), 'What changed from ConstructionBots (Stanford SISL) to this repository',
     20, LIGHT, first=True)
line(tb(s, 0.90, 6.50, 11.5, 0.5), 'Hyunjoon Cha   ·   July 2026', 16, RGBColor(0x9A, 0xA6, 0xBC), first=True)
_page[0] += 1   # 표지엔 foot()를 안 부르므로 페이지 카운터를 직접 1 증가시킴
notes(s, """
[0:00–0:45] Framing.
One sentence: I started from SISL's ConstructionBots — a closed-world TAMP solver that assembles a
model as fast as possible — and ended with a *self-evolving surrogate world model* that lets that
planner survive events it was never designed for, at ~0 planner compute.

The talk has three acts:
  Act 1 (~4 min) — I made the world breakable and gave it a repair vocabulary (OOD + DSL + verifier).
  Act 2 (~3 min) — Two decision-makers over that vocabulary: an LLM layer and a Dec-POMDP/MARL policy.
  Act 3 (~6 min) — The pivot: don't chase sim2real; amortize the expensive twin. E1–E4 results.

VIDEO PLAN for this deck (4 clips, all already runnable in ConstructionBots.jl/tools/):
  V1 · slide 3  — baseline healthy build (tools/demo_wholebuild_anim.jl), 15 s, no OOD.
  V2 · slide 7  — LLM layer repairing a breakdown (tools/demo_respec_replace_anim.jl), 30 s.
  V3 · slide 9  — RL producer on the same fault (tools/demo_rl_replace_anim.jl), 20 s.
  V4 · slide 12 — HEADLINE: surrogate scoring 5 macros in imagination over an OOD stream
                  (tools/demo_surrogate_stream_anim.jl), 45–60 s. This is the one to show if short on time.
Each MeshCat run writes results/tractor/greedy_RVO_Dispersion_TangentBug/visualization.html —
screen-record the browser, trim, save as .mp4, then Insert ▸ Video ▸ This Device into the dashed box.
""")

# ================================================================ 2 · thesis
# [슬라이드 2 · 논지(thesis)] 한 장으로 요약한 프로젝트 arc. "적응이 목표였지만, 진짜 기여는 amortization(비용 분할상환)"이라는 핵심 메시지.
s = slide()
line(tb(s, 0.55, 0.30, 9.0, 0.4), 'THE ARC IN ONE SLIDE', 13, ORANGE, True, first=True)
bar(s)
# 왼쪽에 세로 오렌지 강조 막대(본문 인용 느낌)
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.10), Inches(2.15), Inches(0.12), Inches(3.20))
acc.fill.solid(); acc.fill.fore_color.rgb = ORANGE; acc.line.fill.background(); acc.shadow.inherit = False
tf = tb(s, 1.50, 2.10, 10.6, 3.4)
line(tf, 'ConstructionBots solves a world that never surprises it.', 30, NAVY, True, first=True, space_after=14)
line(tf, 'I made it surprise-able — then discovered the real', 30, INK, True, space_after=0)
line(tf, 'bottleneck was not adaptation. It was the cost of', 30, INK, True, space_after=0)
line(tf, 'asking the planner what to do.', 30, INK, True, space_after=14)
line(tf, 'So the planner became something you can learn.', 30, TEAL, True)
foot(s, 'Adaptation was the goal. Amortization is the contribution.')
notes(s, """
[0:45–1:30] The thesis.
Read the four lines. The pivot is the punchline of the whole talk, so plant it early: the interesting
problem turned out not to be "can we adapt?" (we could, twice over) but "can we afford to?" — one true
re-plan is 70 seconds of MILP + RVO, and a decision over 5 candidate repairs is ~348 s. That is the
number that reframed the project.
""")

# ================================================================ 3 · baseline (VIDEO 1)
# [슬라이드 3 · 출발점] SISL 원본 ConstructionBots가 이미 하던 것(완결형 TAMP 스택). 왼쪽=기능 카드, 오른쪽=데모영상 V1 자리.
s = slide()
head(s, 'THE STARTING POINT', 'What ConstructionBots (SISL) already did')
card(s, 0.70, 2.05, 5.6, 3.55, accent=BLUE)   # 왼쪽 카드: 원본이 가진 기능 목록
tf = tb(s, 1.00, 2.40, 5.0, 3.0)
line(tf, 'A complete TAMP stack', 21, BLUE, True, first=True, space_after=12)
line(tf, '· scene tree + schedule DAG', 15, INK, space_after=7)
line(tf, '· MILP / greedy task assignment', 15, INK, space_after=7)
line(tf, '· transport-team formation', 15, INK, space_after=7)
line(tf, '· RVO2 / TangentBug reactive motion', 15, INK, space_after=7)
line(tf, '· objective: minimize makespan', 15, INK, space_after=12)
line(tf, 'Solve once → execute → done.', 15, GRAY, True, italic=True)   # 원본의 핵심 가정: 한 번 풀고 그대로 실행
# 오른쪽: 방해 없는 정상 빌드 데모영상(V1) 자리
video_ph(s, 6.70, 2.05, 5.9, 3.55, 'V1 — baseline build, no disruption',
         'demo_wholebuild_anim.jl · 15 s')
foot(s, 'Correct, optimal — and completely closed-world. Nothing in it can go wrong.')
notes(s, """
[1:30–2:45] The baseline.
ConstructionBots is genuinely good at what it does: it takes an LDraw model, builds a schedule DAG,
assigns tasks with a MILP (or greedy), forms carrying teams, and drives them with RVO2 collision
avoidance. The objective is a single number: makespan.

The structural assumption is that the world it planned for is the world it executes in. There is no
robot failure, no battery, no keep-out region, no re-plan trigger. The plan is computed once.

>>> VIDEO 1 (15 s, silent, loop): a clean tractor build finishing.
    Run: julia +lts --project=. tools/demo_wholebuild_anim.jl
    Say over it: "this is the whole original story — and it is the only story it can tell."
""")

# ================================================================ 4 · the gap
# [슬라이드 4 · 공백(gap)] 실제 조립 현장은 닫힌세계가 아니다. 3가지 OOD(분포 밖) 사건을 원형 불릿으로 나열.
s = slide(dark=True)
head(s, 'THE GAP', 'Real assembly floors are not closed worlds', dark=True)
# (제목, 설명, 색) 3종 OOD 사건: 로봇고장 / 진입금지구역 / 배터리 방전
items = [
    ('Robot breakdown', 'a carrier dies mid-haul, holding a chain of pending tasks', ORANGE),
    ('Keep-out zone',   'a region becomes unsafe — and it overlaps a staging circle', TEAL),
    ('Battery depletion', 'a robot runs flat; how flat decides what you should do', BLUE),
]
for i, (t, d, c) in enumerate(items):   # enumerate: 인덱스 i로 세로 위치를 계산하며 항목 언패킹
    y = 2.35 + i * 1.30                  # i번째 항목의 세로 y좌표(1.30인치 간격으로 아래로 배치)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(y), Inches(0.5), Inches(0.5))   # 색 원형 불릿
    o.fill.solid(); o.fill.fore_color.rgb = c; o.line.fill.background(); o.shadow.inherit = False
    line(tb(s, 1.65, y - 0.08, 10.5, 0.6), t, 25, WHITE, True, first=True)   # 사건 제목
    line(tb(s, 1.65, y + 0.50, 10.5, 0.5), d, 16, LIGHT, first=True)         # 사건 설명
line(tb(s, 0.85, 6.30, 11.6, 0.8), 'A closed-world plan does not degrade under these. It just fails.',
     20, ORANGE, True, first=True)
foot(s, dark=True)
notes(s, """
[2:45–3:30] The gap I chose to attack.
Three out-of-distribution (OOD) event classes, picked because they hit three *different* parts of the
stack: an agent (breakdown), the space (keep-out zone), a resource (battery). Together they force a
repair vocabulary that is more than one trick.

Key point: in the original code none of these exist — so step one was not "adapt", it was
"make the twin capable of breaking, and make the breakage measurable."
""")

# ================================================================ 5 · layer 1: OOD injection
# [슬라이드 5 · 변화1: 부술 수 있는 twin] 시뮬레이터를 "재현 가능하게 고장낼 수 있게" 만든 5개 신규 파일 목록.
s = slide()
head(s, 'CHANGE 1 · THE BREAKABLE TWIN', 'I made the world able to go wrong — reproducibly')
# (파일명, 역할설명, 색) 새로 추가한 OOD 레이어 구성요소들
rows = [
    ('ood_stream.jl',    'seeded OOD stream — events fire at closed-node thresholds, so a run is reproducible', ORANGE),
    ('ood_injection.jl', 'fault / battery / zone / team-deadlock injectors into the live PlannerEnv', ORANGE),
    ('battery.jl',       'per-robot state-of-charge accounting + stall gate (SoC = 0 ⇒ the robot stops)', BLUE),
    ('ood_truth.jl',     'ground-truth log of what SHOULD have been done — evaluator-only, never seen by the policy', TEAL),
    ('metrics.jl',       'RunMetrics: completion, realized makespan, energy, min-SoC, reaction & recovery time', TEAL),
]
for i, (f, d, c) in enumerate(rows):   # 각 파일을 한 행씩 그림(왼쪽 색 막대 + 파일명 + 설명)
    y = 2.10 + i * 0.92                 # i번째 행의 세로 위치(0.92인치 간격)
    m = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.80), Inches(y + 0.05), Inches(0.13), Inches(0.62))  # 왼쪽 색 표시 막대
    m.fill.solid(); m.fill.fore_color.rgb = c; m.line.fill.background(); m.shadow.inherit = False
    line(tb(s, 1.10, y, 4.0, 0.5), f, 19, NAVY, True, first=True)      # 파일명(굵게)
    line(tb(s, 4.60, y + 0.06, 8.1, 0.6), d, 15, GRAY, first=True)     # 역할 설명
foot(s, 'The simulator core stays a black box — everything hangs off Ref-registry hooks.')
notes(s, """
[3:30–4:30] Change 1 — the OOD layer (src/navigator/).
Five new pieces. The design constraint I set for myself: do NOT fork the simulator. Every addition
attaches through hook registries (pre_sim_hook, SOC_SPEED_HOOK, BATTERY_STEP_HOOK), so the original
ConstructionBots step loop is untouched and still passes its own tests.

Two things worth calling out because they later saved the whole evaluation:
  · The OOD stream fires at *closed-node* thresholds, not wall-clock steps → byte-reproducible across
    machines, which is what makes paired A/B comparison legitimate.
  · ood_truth.jl records the correct answer at injection time. The controller never sees it; only the
    scorer does. Without this you cannot score a decision, only an outcome.
""")

# ================================================================ 6 · layer 2: DSL + verifier
# [슬라이드 6 · 변화2: 수리 어휘(DSL)] 위 줄=5개 매크로 카드, 아래 줄=처리 파이프라인 체인(사건→DSL→검증→재계획→실행).
s = slide()
head(s, 'CHANGE 2 · THE REPAIR VOCABULARY', 'A closed DSL — and a verifier that owns safety')
# 닫힌(closed) DSL의 전부 = 5개 매크로. 어떤 제안자든 이 5개만 뱉을 수 있다. (매크로명, 짧은설명, 색)
macros = [
    ('NOOP', 'absorb it', GRAY),
    ('Replace', 'hot-swap a spare', ORANGE),
    ('Deprioritize', 'offload heavy hauls', BLUE),
    ('ForbidZone', 'restage around it', TEAL),
    ('ReformTeam', 'rebuild the carrier', NAVY),
]
for i, (m, d, c) in enumerate(macros):   # 위쪽 줄: 5개 매크로를 가로로 나란히 카드로 배치
    x = 0.70 + i * 2.47                   # i번째 카드의 가로 위치(2.47인치 간격)
    card(s, x, 2.05, 2.25, 1.55, accent=c)
    line(tb(s, x + 0.15, 2.42, 1.95, 0.5), m, 15, c, True, first=True, align=PP_ALIGN.CENTER)   # 매크로 이름
    line(tb(s, x + 0.15, 2.95, 1.95, 0.6), d, 12, GRAY, first=True, align=PP_ALIGN.CENTER)      # 한 줄 설명
# 아래쪽 줄: 수리 처리 파이프라인 단계들. verifier(검증기)가 안전을 책임지는 관문(gate) 역할.
chain = [('OOD event', GRAY), ('DSL spec', ORANGE), ('verifier gate', RED), ('MILP re-plan', BLUE), ('execute', TEAL)]
for i, (t, c) in enumerate(chain):
    x = 0.70 + i * 2.47
    card(s, x, 4.30, 2.10, 1.05, accent=None, fill=CARD_L)
    line(tb(s, x + 0.05, 4.62, 2.00, 0.5), t, 15, c, True, first=True, align=PP_ALIGN.CENTER)
    if i < 4:                              # 마지막 단계 뒤에는 화살표를 안 그림(단계 사이에만)
        arrow(s, x + 2.14, 4.72, 0.28, 0.20)
line(tb(s, 0.70, 5.70, 11.9, 0.5), 'admit / reject / repair  —  a proposal that does not verify never reaches the solver',
     15, RED, True, first=True)
foot(s, 'The proposer says WHAT. The solver decides HOW. The verifier decides IF.')
notes(s, """
[4:30–5:30] Change 2 — the DSL + the enactment pipeline (src/respec/).
Five macros is the entire adaptation vocabulary. It is a *closed* union on purpose: whatever proposes
a repair — an LLM, an RL policy, a surrogate, a random baseline — can only emit these, and every one
of them passes the identical verifier → compiler → MILP-replan → resume path.

That closure is what makes the rest of the project possible: it turns "adapt to an OOD" into a
5-way discrete decision, which is (a) safely enactable and (b) learnable.

Note NOOP is a first-class macro, not the absence of a decision. Deciding *not* to intervene is often
the right answer — and that turned out to be the hardest thing to teach (slide 13).
""")

# ================================================================ 6b · by the numbers
# [슬라이드 6b · 변경 규모(diff)] 위=핵심 숫자 카드 4개, 아래=하위 폴더별 코드 규모 표. "원본 core는 건드리지 않았다(additive)"가 요지.
s = slide()
head(s, 'THE DIFF', 'Everything is additive — the SISL core is untouched')
# 상단 큰 숫자 카드: (숫자, 라벨, 색). 강조 포인트=원본 시뮬레이터 fork 파일 0개, hook 4개뿐.
stats = [
    ('~17,000', 'lines added', ORANGE),
    ('104', 'new files', BLUE),
    ('4', 'hooks into the original', TEAL),
    ('0', 'forked simulator files', NAVY),
]
for i, (n, l, c) in enumerate(stats):   # 숫자 카드 4개를 가로로 배치
    x = 0.75 + i * 3.10                  # i번째 카드 가로 위치(3.10인치 간격)
    card(s, x, 2.10, 2.75, 1.85, accent=c)
    line(tb(s, x + 0.15, 2.45, 2.45, 0.8), n, 36, c, True, first=True, align=PP_ALIGN.CENTER)   # 큰 숫자
    line(tb(s, x + 0.15, 3.35, 2.45, 0.5), l, 14, GRAY, first=True, align=PP_ALIGN.CENTER)      # 라벨
# 하단 표: (폴더/대상, 코드량, 내용, 색). 각 하위 시스템이 무엇을 담당하는지 한 줄 요약.
rows = [
    ('src/respec/', '~5,200 LOC', 'DSL · verifier · compiler · re-plan · LLM bridge · OOD injection · replace · restage', ORANGE),
    ('src/navigator/', '~2,050 LOC', 'metrics · baselines B0–B5 · controllers · battery · OOD stream · ground truth', TEAL),
    ('decpomdp/', '~5,200 LOC', 'Dec-POMDP wrapper · macro action head · MAPPO / CTDE trainers', BLUE),
    ('wm4spacecraft/', 'E1–E4', 'oracle dump seam · surrogate export · frontier / drift / ablation analysis', NAVY),
    ('the original src/*.jl', '4 hooks', 'pre_sim_hook · BATTERY_STEP_HOOK · SOC_SPEED_HOOK · ood_inject/respec step seam', GRAY),
]
for i, (f, n, d, c) in enumerate(rows):   # 표 한 행씩: 폴더명 / 코드량 / 설명 3열로 그림
    y = 4.15 + i * 0.48                    # i번째 행 세로 위치(0.48인치 간격, 촘촘한 표)
    m = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.80), Inches(y + 0.03), Inches(0.11), Inches(0.30))  # 왼쪽 색 표식
    m.fill.solid(); m.fill.fore_color.rgb = c; m.line.fill.background(); m.shadow.inherit = False
    line(tb(s, 1.05, y, 2.7, 0.4), f, 14, NAVY, True, first=True)      # 1열: 폴더/대상
    line(tb(s, 3.60, y + 0.02, 1.3, 0.4), n, 13, c, True, first=True)  # 2열: 코드량
    line(tb(s, 5.05, y + 0.02, 7.6, 0.4), d, 13, GRAY, first=True)     # 3열: 내용
foot(s, 'Disable the hooks and upstream ConstructionBots behaves exactly as before.')
notes(s, """
[optional — cut this slide first if you are over time; it is engineering credibility, not argument.]

The scale of the fork: ~17k LOC across two new subsystems in ConstructionBots.jl plus a standalone
decpomdp package and the wm4spacecraft experiment harness. But the number I actually care about is the
last one: FOUR hook points touch the original simulator, all inert by default.

That was a deliberate constraint. If the OOD layer had required forking the solver, no claim about the
solver's behavior would carry over, and nothing I built could be re-fitted onto someone else's twin.
Additive-only is what makes the "the recipe transfers" claim on the limits slide honest.
""")

# ================================================================ 7 · producer A: LLM (VIDEO 2)
# [슬라이드 7 · 변화3a: 제안자 A = LLM] 자연어 사건 → DSL 매크로로 바꾸는 LLM 재명세(respec) 레이어. 오른쪽=데모영상 V2.
s = slide(dark=True)
head(s, 'CHANGE 3a · PRODUCER A', 'The LLM re-specification layer', dark=True)
tf = tb(s, 0.75, 2.05, 5.9, 3.6)   # 왼쪽 설명 글틀(예시: 자연어 → ReplaceAgent DSL 변환 흐름)
line(tf, 'NL event  →  DSL spec', 22, ORANGE, True, first=True, space_after=14)
line(tf, '"Robot R3 is immobile in the south bay"', 16, WHITE, italic=True, space_after=8)
line(tf, '↓   tool-pinned, schema-validated call', 14, LIGHT, space_after=8)
line(tf, 'ReplaceAgent(R3)', 18, ORANGE, True, space_after=16)
line(tf, '· zero-shot — no training, no OOD distribution assumed', 14, LIGHT, space_after=7)
line(tf, '· exact-id grounding, no fuzzy matching', 14, LIGHT, space_after=7)
line(tf, '· the verifier, not the LLM, carries the guarantee', 14, LIGHT, space_after=7)
line(tf, '· cost: seconds of API latency per decision', 14, RED, space_after=0)   # 단점: 결정마다 API 지연(초 단위) 비용
# 오른쪽: LLM이 로봇 고장을 수리하는 데모영상(V2) 자리
video_ph(s, 7.00, 2.05, 5.6, 3.6, 'V2 — LLM repairs a breakdown',
         'demo_respec_replace_anim.jl · 30 s')
foot(s, 'Open-world reasoning, bounded by a closed grammar.', dark=True)
notes(s, """
[5:30–6:30] Producer A — the LLM layer (src/respec/llm_service/).
This was the first adaptive controller. The LLM reads the OOD event as natural language and emits one
DSL macro through a pinned tool call, so it structurally cannot invent an unsafe action. It is
zero-shot: it never trains on the OOD distribution, which is exactly the property you want for events
you did not anticipate.

>>> VIDEO 2 (30 s): tools/demo_respec_replace_anim.jl
    Show: fault fires mid-build → the sidebar prints the emitted ReplaceAgent spec → a spare drives out
    of the depot, hot-swaps into the dead robot's identity, and the build continues to completion.
    Say over it: "the LLM only said WHAT. The MILP re-solved HOW."

    (Alternate if you prefer the spatial case: demo_respec_forbidzone_anim.jl — a zone lands on a
     staging circle and ForbidZone restages the assembly. Only show one; time is tight.)
""")

# ================================================================ 8 · producer B: MARL
# [슬라이드 8 · 변화3b: 제안자 B = MARL] 같은 DSL 어휘를 쓰는 학습형 결정자(Dec-POMDP). 왼쪽=설계 카드, 오른쪽=솔직한 실패담(붕괴→재설계) 카드.
s = slide()
head(s, 'CHANGE 3b · PRODUCER B', 'The Dec-POMDP / MARL policy — and what it taught me')
card(s, 0.70, 2.05, 5.85, 2.35, accent=BLUE)   # 왼쪽 카드: MARL 설계(매크로 head, 관측차원, MAPPO/CTDE)
tf = tb(s, 1.00, 2.40, 5.25, 1.9)
line(tf, 'A learned decision-maker over the SAME vocabulary', 17, BLUE, True, first=True, space_after=10)
line(tf, '· macro head lifted to the DSL repertoire (not throttle)', 14, INK, space_after=6)
line(tf, '· observation + OOD-signal block (obs dim 50 → 63)', 14, INK, space_after=6)
line(tf, '· MAPPO / CTDE, held-out severity split for fairness', 14, INK, space_after=0)
card(s, 6.80, 2.05, 5.8, 2.35, accent=RED)   # 오른쪽 카드(빨강): 학습 붕괴 → PPO 최적화 문제로 진단 → 이벤트형 매크로 MDP로 재설계
tf = tb(s, 7.10, 2.40, 5.2, 1.9)
line(tf, 'The honest result', 17, RED, True, first=True, space_after=10)
line(tf, 'The per-step throttle MDP collapsed. An ablation ruled out', 14, INK, space_after=6)
line(tf, 'the reward (it is monotone in throttle) — it was a PPO', 14, INK, space_after=6)
line(tf, 'optimization failure. So I redesigned it as an', 14, INK, space_after=6)
line(tf, 'event-triggered macro MDP, symmetric with the LLM.', 14, INK, True, space_after=0)
tf = tb(s, 0.70, 4.85, 11.9, 1.4)
line(tf, 'Why this matters for the arc:', 19, NAVY, True, first=True, space_after=8)
line(tf, 'Two very different producers — zero-shot reasoning and a trained policy — end up making the same '
         'kind of call: pick one of five macros. That symmetry is what later let me swap in a third producer.',
     17, GRAY, space_after=0)
foot(s)
notes(s, """
[6:30–7:30] Producer B — the Dec-POMDP (decpomdp/).
I built the MARL arm to answer "does an LLM actually beat a learned policy at this?" To make that a
fair question I had to lift the action space: the original wrapper only let agents throttle (HOLD/SLOW/GO),
which can bias motion but structurally cannot recover a dead robot. So MARL got the same five macros.

Be honest here — it is the strongest part of the story: the first formulation (per-step throttle) collapsed
in training. I ran an ablation that showed reward was monotone in the action, so the reward was fine; the
failure was PPO's, on a horrible per-step credit-assignment problem. Rather than tune it forever I killed
the formulation and re-cast it as an event-triggered decision — one macro choice per OOD event — which is
*exactly* the LLM's interface.

That is the accidental discovery that set up the pivot: the decision problem is small and discrete.
Small discrete decisions are cheap to LEARN — but each one still costs a 70-second true-planner call to
EVALUATE. That asymmetry is the whole rest of the talk.
""")

# ================================================================ 9 · the pivot
# [슬라이드 9 · 전환점(pivot)] 핵심 슬라이드. 왼쪽=피한 함정(sim2real), 오른쪽=재구성(amortization=값비싼 twin을 분할상환). 아래=비용 숫자.
s = slide(dark=True)
head(s, 'THE PIVOT', 'Adapting inside a toy world is not the deliverable', dark=True)
card(s, 0.75, 2.10, 5.7, 2.15, dark=True, accent=RED)   # 왼쪽 카드(빨강): 거부한 함정 = 장난감 세계 안에서 튜닝(sim2real gap)
tf = tb(s, 1.05, 2.45, 5.1, 1.7)
line(tf, 'The trap I refused', 16, RED, True, first=True, space_after=9)
line(tf, 'ConstructionBots is a simplified world. Tuning an adaptive', 14, WHITE, space_after=5)
line(tf, 'controller inside it does not transfer to a real assembly', 14, WHITE, space_after=5)
line(tf, 'floor — the sim2real gap eats the result.', 14, WHITE, space_after=0)
card(s, 6.90, 2.10, 5.7, 2.15, dark=True, accent=TEAL)   # 오른쪽 카드(틸): 재구성 = 실제 우주선 제조도 비싼 digital twin 위에서 계획한다 → sim2real이 아니라 amortization
tf = tb(s, 7.20, 2.45, 5.1, 1.7)
line(tf, 'The reframe', 16, TEAL, True, first=True, space_after=9)
line(tf, 'Real spacecraft manufacturing does not plan from scarce real', 14, WHITE, space_after=5)
line(tf, 'data either — it plans over an expensive DIGITAL TWIN.', 14, WHITE, space_after=5)
line(tf, 'So the deployable move is not sim2real. It is amortization.', 14, WHITE, True, space_after=0)
tf = tb(s, 0.75, 4.75, 11.85, 1.6)
line(tf, 'The expensive thing is the query, not the world.', 26, ORANGE, True, first=True, space_after=10)
line(tf, 'One true-planner label = 69.5 s (MILP + full RVO sim).   One 5-macro decision = ~348 s.',
     19, WHITE, True, space_after=6)
line(tf, 'Learn a fast surrogate of the twin · decide in imagination · verify sparingly.  (MuZero-style value equivalence)',
     16, LIGHT, space_after=0)
foot(s, dark=True)
notes(s, """
[7:30–8:45] THE PIVOT — the most important slide. Slow down here.

The argument, in the order I want it heard:
1. Yes, I can make ConstructionBots adapt — twice over (LLM, MARL). But ConstructionBots is a toy.
   Adaptive control tuned in a toy does not survive contact with a real facility.
2. So do we need a higher-fidelity sim? That is the sim2real treadmill, and it never ends.
3. The reframe: nobody plans a spacecraft assembly line from real data — the data does not exist.
   They plan over a high-fidelity digital twin. The twin IS the world model. It is not inaccurate;
   it is just EXPENSIVE.
4. Therefore the deployable contribution is not closing a fidelity gap. It is making the expensive
   known world cheap to query — MuZero's move, applied to a TAMP planner instead of Atari.

Numbers to land: each candidate repair must be *simulated* to know if it was right — 69.5 s. Five
candidates per OOD = 348 s per decision. That is why re-planning under disruption is unaffordable,
and that is what a surrogate removes.

Precision worth stating out loud: the surrogate does NOT imitate the simulator's dynamics. It maps
(state, OOD event, candidate macro) → outcome metrics. It is a value-equivalent decision surrogate.
""")

# ================================================================ 10 · new architecture
# [슬라이드 10 · 변화4: 새 아키텍처] 좌→우 흐름: OOD사건 → Producer(제안) → Surrogate(상상 채점) → top-k만 검증 → 실행. 아래 점선박스=자기진화 재학습 루프.
s = slide()
head(s, 'CHANGE 4 · THE NEW ARCHITECTURE', 'Producer → surrogate → sparse verification')
# (제목, 부제, 색, x좌표) 5단계 파이프라인 박스. surrogate가 5개 매크로를 planner 호출 없이 "상상"으로 채점.
boxes = [
    ('OOD event',        'fault / zone / battery',        GRAY,   0.70),
    ('Producer',         'LLM  ·  RL  ·  all-macros',     ORANGE, 3.20),
    ('Surrogate',        'scores 5 macros in imagination', TEAL,  5.70),
    ('Verify top-k',     'true planner, k ≪ 5',           BLUE,   8.20),
    ('Enact',            'MILP re-plan → execute',        NAVY,  10.70),
]
for t, d, c, x in boxes:   # 여기선 x좌표를 튜플에 직접 담아 언패킹(간격이 균일하지 않아서)
    card(s, x, 2.35, 2.15, 1.60, accent=c)
    line(tb(s, x + 0.07, 2.70, 2.01, 0.5), t, 15, c, True, first=True, align=PP_ALIGN.CENTER)   # 단계 제목
    line(tb(s, x + 0.07, 3.22, 2.01, 0.6), d, 12, GRAY, first=True, align=PP_ALIGN.CENTER)      # 단계 부제
    if x < 10.5:                                   # 마지막(Enact, x=10.70) 뒤엔 화살표 생략
        arrow(s, x + 2.20, 3.05, 0.28, 0.20)
# feedback loop  / 자기진화 루프: 예측 오차 감시(drift 감지) → 실제 planner로 재라벨 → surrogate 재학습
fb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.20), Inches(4.55), Inches(9.80), Inches(1.05))
fb.fill.solid(); fb.fill.fore_color.rgb = RGBColor(0xE3, 0xEE, 0xEC)
fb.line.color.rgb = TEAL; fb.line.width = Pt(1.25); fb.line.dash_style = 4; fb.shadow.inherit = False
tf = tb(s, 3.45, 4.72, 9.3, 0.85)
line(tf, 'THE SELF-EVOLVING LOOP', 13, TEAL, True, first=True, space_after=4)
line(tf, 'residual monitor (Page-Hinkley/CUSUM) → drift detected → relabel with the true planner → retrain the surrogate',
     15, INK, space_after=0)
foot(s, 'The producer proposes broadly. The surrogate ranks for free. The planner is called only when it must be.')
notes(s, """
[8:45–9:45] Change 4 — the architecture the repo now implements.
Read it left to right: an OOD fires; a producer proposes candidate macros (the LLM proposing is one
option — RL is another, "propose all five" is the ceiling); the learned surrogate predicts each
candidate's *outcome* without running the planner; only the top-k are actually verified; the best
verified one is enacted through the same verifier→MILP path as before.

The dashed loop is what makes it self-evolving rather than a static regressor: the surrogate's
prediction error on the one candidate it did verify is a free residual signal. A CUSUM monitor watches
it; when the OOD distribution drifts, it escalates verification, relabels, and retrains. Without this
loop the cheap surrogate is one distribution shift away from catastrophe — which I show on slide 12.

Note the LLM was not thrown away in the pivot. It got demoted to one interchangeable producer among
several — which is exactly the modularity claim from the earlier deck ("the LLM is a swappable part").
""")

# ================================================================ 11 · E1 results
# [슬라이드 11 · 결과 E1] 성능표. surrogate가 planner 호출 0번으로 oracle 수준 결정(regret 0). 각 방법의 후회도/치명적 오선택/호출수 비교.
s = slide()
head(s, 'RESULT · E1', 'The surrogate makes the planner\'s decision at zero planner calls')
hdr = ['method', 'decision regret', 'catastrophic picks', 'planner calls / decision']   # 표 헤더 4개 열
# (방법, decision regret, 치명적 오선택 수, planner 호출/결정, 색). 우리 것(surrogate)만 호출 0으로 regret 0.
data = [
    ('learned surrogate (ours)', '0.000', '0 / 32', '0', TEAL),
    ('best state-blind policy',  '0.281', '9 / 32', '5', GRAY),
    ('random',                   '0.844', '15 / 32', '5', GRAY),
    ('true-planner oracle',      '0.000', '0 / 32', '5  (≈348 s)', BLUE),
]
xs = [0.80, 5.20, 7.90, 10.40]   # 4개 열의 가로 시작 x좌표
ws = [4.30, 2.60, 2.40, 2.20]    # 4개 열의 폭
for j, h in enumerate(hdr):      # 헤더 행: h.upper()로 대문자화해 오렌지색으로 출력
    line(tb(s, xs[j], 2.10, ws[j], 0.4), h.upper(), 12, ORANGE, True, first=True)
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.80), Inches(2.52), Inches(11.8), Inches(0.02))  # 헤더 아래 구분선
sep.fill.solid(); sep.fill.fore_color.rgb = LIGHT; sep.line.fill.background(); sep.shadow.inherit = False
for i, (m, r, c, k, col) in enumerate(data):   # 데이터 행들
    y = 2.70 + i * 0.62                          # i번째 행 세로 위치
    bold = (i == 0)                              # 첫 행(=우리 방법)만 굵게+색으로 강조
    line(tb(s, xs[0], y, ws[0], 0.5), m, 16, col if bold else INK, bold, first=True)   # 1열: 방법 이름
    for j, v in enumerate([r, c, k], start=1):   # 나머지 3열(값들)을 j=1..3 위치에 출력
        line(tb(s, xs[j], y, ws[j], 0.5), v, 16, col if bold else INK, bold, first=True)
card(s, 0.80, 5.35, 5.75, 1.05, fill=CARD_L)   # 하단 왼쪽 카드: 예측 정확도(R²/MAE)
line(tb(s, 1.05, 5.55, 5.3, 0.7), 'Held-out outcome prediction:  R² = 0.83,  MAE 13 nodes',
     15, NAVY, True, first=True)
card(s, 6.85, 5.35, 5.75, 1.05, fill=CARD_L)   # 하단 오른쪽 카드: 이득의 통계적 유의성(paired bootstrap 95% CI가 0을 제외)
line(tb(s, 7.10, 5.55, 5.3, 0.7), 'Paired bootstrap 95% CI on the gain:  +0.281 [+0.13, +0.44]',
     15, NAVY, True, first=True)
foot(s, '100% planner-compute reduction at oracle-quality decisions  ·  32 admissible instances, leave-one-out')
notes(s, """
[9:45–10:45] E1 — planner amortization. The headline number.
Decision regret = true cost of the macro you picked minus true cost of the oracle-best macro,
normalized per scenario. Zero regret means: the surrogate chose exactly what a 348-second brute-force
sweep of all five macros would have chosen — at zero planner calls.

Two guards against over-claiming, both on the slide:
  · R² 0.83 says it is genuinely predicting the outcome, not memorizing a lookup — but note the ranking
    is perfect *despite* ~13 nodes of prediction error. That is the value-equivalence point: wrong
    values, right decision.
  · The paired bootstrap CI excludes zero, so beating the state-blind baseline is real, not noise.

The "catastrophic picks" column is the one to point at: the state-blind policy picks Replace on every
zone event — and Replace leaves the build INCOMPLETE there. Nine times out of 32.
""")

# ================================================================ 12 · E2/E3/E4 (VIDEO 4)
# [슬라이드 12 · 결과 E2/E3/E4] 왼쪽=세 실험 요약(싸다/drift에 안전/Pareto front), 오른쪽=헤드라인 데모영상 V4(surrogate 실시간 결정).
s = slide(dark=True)
head(s, 'RESULTS · E2 · E3 · E4', 'Cheap, safe under drift, and on the Pareto front', dark=True)
# (실험라벨, 한줄결과, 색). E2=더 싸다, E3=새 OOD에도 재학습으로 회복, E4=full 조합만 Pareto front.
res = [
    ('E2', 'LLM → surrogate matches LLM → exact solver at 67% fewer planner calls', ORANGE),
    ('E3', 'Under a novel OOD type the frozen surrogate collapses (regret 1.00); the active retrain loop '
           'recovers in one step (0.11) at 54% fewer calls than the oracle', TEAL),
    ('E4', 'Ablation grid: only the full cell (producer + surrogate + retraining) is both cheap and '
           'drift-robust — it sits on the quality–compute Pareto front', BLUE),
]
for i, (k, d, c) in enumerate(res):   # 실험 3개를 위→아래로 배치(색 배지 + 설명)
    y = 2.15 + i * 1.15                # i번째 항목 세로 위치
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.80), Inches(y), Inches(0.85), Inches(0.72))  # E2/E3/E4 색 배지
    b.fill.solid(); b.fill.fore_color.rgb = c; b.line.fill.background(); b.shadow.inherit = False
    line(tb(s, 0.80, y + 0.13, 0.85, 0.5), k, 18, WHITE, True, first=True, align=PP_ALIGN.CENTER)   # 배지 안 라벨
    line(tb(s, 1.95, y - 0.02, 5.0, 1.1), d, 15, WHITE, first=True)                                 # 오른쪽 설명
# 오른쪽: 헤드라인 데모영상(V4) 자리 — surrogate가 5개 매크로를 실시간으로 상상 채점하는 장면
video_ph(s, 7.20, 2.05, 5.4, 3.55, 'V4 — the surrogate deciding, live',
         'demo_surrogate_stream_anim.jl · 45–60 s')
line(tb(s, 0.80, 5.95, 6.2, 0.9), 'Drop the retraining loop and the cheap surrogate becomes unsafe. '
     'That loop is the "self-evolving" part.', 15, LIGHT, True, first=True)
foot(s, dark=True)
notes(s, """
[10:45–12:15] E2/E3/E4 — do not read the table, tell the story of each letter.
  E2: the prior-art regime is "LLM proposes → exact solver verifies every candidate" (arXiv 2506.18178).
      Same decisions, one third of the calls, because the surrogate's top-1 is already the solver's pick.
  E3: I drift the stream — phase A is robot faults, then a novel blocking-zone type appears. A frozen
      surrogate keeps answering "Replace" and is catastrophically wrong forever. The active loop detects
      the residual spike, relabels, retrains, and is back to ~0 regret one step later.
  E4: the ablation grid IS the related-work table — each cell is one prior-art regime. Only the full
      system is on the Pareto front.

>>> VIDEO 4 — THE HEADLINE DEMO (45–60 s): tools/demo_surrogate_stream_anim.jl
    A stream of 4 different OOD events hits one build. At each, the docked sidebar shows the surrogate
    scoring ALL FIVE macros in imagination — no planner call — and the argmax being enacted, with a live
    per-robot battery SoC bar draining against the playhead.
    Narrate: "every number you see appear in that sidebar would have cost 70 seconds of MILP + RVO to
    obtain honestly. The surrogate produces all five in microseconds — and picks the same one."
    If the talk is running long, THIS is the clip to keep and the others to cut.
""")

# ================================================================ 13 · rigor / what broke
# [슬라이드 13 · 실제로 한 일(엄밀성)] 작업 대부분은 "벤치마크를 정직하게 만드는" 디버깅이었다. 발견한 harness 버그/방법론 오류 4가지.
s = slide()
head(s, 'WHAT THE WORK ACTUALLY WAS', 'Most of it was making the benchmark honest')
# (문제 제목, 상세 설명, 색). RED=진짜 버그, ORANGE/TEAL=방법론·모델 선택 교정.
rows = [
    ('The battery OOD was hitting a parked spare.',
     'The injector picked argmax(SoC) — which, once batteries are on, is always an idle spare that owns no '
     'work. Every macro produced a byte-identical sim, so "battery is harmless" was a bug, not a finding.', RED),
    ('The NOOP arm was not really NOOP.',
     'A stalled robot raised a breakdown alarm that the background policy silently answered with a Replace — '
     'inside the do-nothing arm. The studied decision now owns its own consequences.', RED),
    ('Restraint had been deleted from the dataset.',
     'The admissibility filter kept only instances where doing nothing was wrong, so the model could never '
     'learn to do nothing. Cost-aware labels (y = closed − λ·cost) let NOOP win when the OOD is absorbed anyway.', ORANGE),
    ('A linear surrogate cannot do this job.',
     'The right macro flips at a threshold (SoC ≈ 0.15, zone overlap ≈ 0.4). An additive model gives every '
     'macro the same slope, so the curves never cross. Trees represent thresholds natively.', TEAL),
]
for i, (t, d, c) in enumerate(rows):   # 문제 4개를 위→아래로: 왼쪽 색 막대 + 제목 + 상세
    y = 1.98 + i * 1.14                 # i번째 항목 세로 위치
    m = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.80), Inches(y + 0.02), Inches(0.13), Inches(0.88))  # 왼쪽 색 막대
    m.fill.solid(); m.fill.fore_color.rgb = c; m.line.fill.background(); m.shadow.inherit = False
    line(tb(s, 1.10, y, 11.4, 0.42), t, 18, NAVY, True, first=True)         # 문제 제목
    line(tb(s, 1.10, y + 0.42, 11.4, 0.72), d, 13, GRAY, first=True)        # 상세 설명
foot(s, 'Same OOD kind, opposite correct answers — H(best | kind) went from 0 bits to 0.95.')
notes(s, """
[12:15–13:15] The unglamorous half — and the part a research audience should trust most.
E1 initially looked *too* good, and the reason was that the task was trivially easy: the OOD kind alone
determined the best macro (fault→Replace, zone→ForbidZone), zero bits of residual entropy. A one-hot on
the event type would score zero regret. That is not a world model; that is a lookup table.

Chasing that down surfaced two genuine harness bugs (the top two rows) that had been silently making
battery events consequence-free, plus one methodological error of my own (the admissibility filter had
deleted every instance where restraint was correct).

After fixing them the same OOD kind gets opposite correct answers depending on severity:
  SoC 0.05 → Replace · SoC 0.20 → NOOP (replacing is WORSE) · SoC 0.5 → NOOP
  zone overlap 0.31 → NOOP (the restage costs more than it saves) · overlap 0.49 → ForbidZone
Now the best fixed macro *per OOD kind* — the strongest state-blind baseline that exists — is wrong on
roughly half the instances. That is a real decision problem, and it is the one the demo shows.
""")

# ================================================================ 14 · limits
# [슬라이드 14 · 정직한 한계] 아직 주장하지 않을 것 3가지(twin은 대리물, mid-build Replace 미완주, surrogate는 아직 GNN 아님). 번호 원 + 제목 + 설명.
s = slide()
head(s, 'HONEST LIMITS', 'What I would not yet claim')
lims = [   # (한계 제목, 상세 설명) 3개
    ('The twin is a proxy, not a spacecraft.',
     'Results run on the ConstructionBots tractor model. The transferable deliverable is the method — DSL, '
     'OOD taxonomy, verify-top-k, retrain loop — re-fittable onto a facility\'s own high-fidelity twin. '
     'The weights do not transfer; the recipe does.'),
    ('A mid-build Replace still does not finish the build.',
     'Diagnosed to one carrier that never reaches its deposit, not to broad congestion. The DECISION is '
     'right (it closes more nodes than doing nothing); the run ends incomplete. Fix implemented, off by '
     'default — turning it on changes the twin\'s dynamics and invalidates every oracle label.'),
    ('The surrogate is a feature model, not yet a GNN.',
     'The milestone gate in the proposal has been met, so a message-passing GNN over (scene tree + schedule '
     'DAG) is the validated next investment — not a blocker.'),
]
for i, (t, d) in enumerate(lims):   # 한계 3개를 위→아래로: 번호 원 + 제목 + 설명
    y = 2.10 + i * 1.55              # i번째 항목 세로 위치(1.55인치 간격)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.80), Inches(y + 0.02), Inches(0.42), Inches(0.42))  # 오렌지 번호 원
    o.fill.solid(); o.fill.fore_color.rgb = ORANGE; o.line.fill.background(); o.shadow.inherit = False
    line(tb(s, 0.80, y + 0.06, 0.42, 0.4), str(i + 1), 15, WHITE, True, first=True, align=PP_ALIGN.CENTER)  # 원 안 번호(1,2,3)
    line(tb(s, 1.50, y, 11.1, 0.45), t, 19, NAVY, True, first=True)         # 한계 제목
    line(tb(s, 1.50, y + 0.46, 11.1, 1.0), d, 13, GRAY, first=True)         # 상세 설명
foot(s, 'Each of these is logged by the pipeline\'s own self-checks — they are the next chapter, not a hidden failure.')
notes(s, """
[13:15–14:00] Limits — say these before anyone asks; it buys credibility for everything above.
The one that will be challenged hardest is #2, so pre-empt it: the Replace decision is *correct* by the
metric (it closes 172 nodes vs 151 for doing nothing) but the run still ends incomplete. I traced it with
a wedge debugger to a single formed carrier that never arrives at its deposit — the waiting teams snap
into their carrying slots, become obstacles in its path, and it never arrives. That is a simulator
dynamics bug, not a decision bug, and the fix exists behind a flag. I have not enabled it because doing
so changes the twin, which would force regenerating every oracle label the surrogate was trained on.
""")

# ================================================================ 15 · summary
# [슬라이드 15 · 요약] 진화 4단계를 카드로 나열(SISL기반 → +OOD&DSL → +제안자 → +surrogate 루프) + 마지막 핵심 메시지.
s = slide(dark=True)
head(s, 'SUMMARY', 'What changed, and why it matters', dark=True)
# (단계명, 2줄설명, 색). 설명 안 '\n'은 카드 안에서 줄바꿈 위치. 단계가 오른쪽으로 갈수록 발전.
stages = [
    ('SISL baseline', 'closed-world TAMP\nsolve once, execute', GRAY),
    ('+ OOD & DSL', 'a breakable twin +\na verified repair vocabulary', ORANGE),
    ('+ Producers', 'LLM re-spec  ·  Dec-POMDP\nsame 5-macro decision', BLUE),
    ('+ Surrogate loop', 'decide in imagination,\nverify sparingly, retrain on drift', TEAL),
]
for i, (t, d, c) in enumerate(stages):   # 4단계 카드를 가로로 배치하고 사이에 화살표
    x = 0.75 + i * 3.10                   # i번째 카드 가로 위치
    card(s, x, 2.30, 2.75, 2.30, dark=True, accent=c)
    line(tb(s, x + 0.18, 2.70, 2.4, 0.6), t, 16, c, True, first=True, align=PP_ALIGN.CENTER)   # 단계 제목
    tfx = tb(s, x + 0.18, 3.35, 2.4, 1.1)
    for j, ln in enumerate(d.split('\n')):   # 설명을 '\n' 기준으로 잘라 여러 줄로 출력(j==0이 첫 문단)
        line(tfx, ln, 13, LIGHT, first=(j == 0), align=PP_ALIGN.CENTER, space_after=4)
    if i < 3:                                # 마지막 카드 뒤엔 화살표 생략
        arrow(s, x + 2.83, 3.35, 0.24, 0.18, ORANGE)
tf = tb(s, 0.75, 5.20, 11.85, 1.5)
line(tf, 'The contribution is not "an LLM adapts a robot fleet".', 22, WHITE, True, first=True, space_after=8)
line(tf, 'It is that an expensive TAMP planner can be amortized into a learned world model that keeps '
         'making the planner\'s decisions — and repairs itself when the world moves.', 19, ORANGE, True, space_after=0)
foot(s, dark=True)
notes(s, """
[14:00–15:00] Close.
Walk the four boxes left to right in one breath, then land the last two lines.

If you take one thing: the original ConstructionBots asks a solver "what is the fastest plan?" once.
This repo asks a *learned model of that solver* "what should I do now?" every time the world breaks —
and only pays the solver when the model admits it is unsure.

Q&A ammunition:
  · "Why not just make the sim faster?" — RVO + MILP is the fidelity; making it cheap by simplifying it
    is exactly the sim2real mistake. Amortize, don't degrade.
  · "Why not end-to-end RL?" — tried it (slide 8). And an RL world model learns the *dynamics*; I only
    need the *decision*, which is a far smaller and better-posed target (value equivalence).
  · "Is the LLM necessary?" — E4 cell B says no for quality, yes for compute: it narrows the candidate
    set, which is where the calls are saved.
""")

# ================================================================ 16 · demo appendix
# [슬라이드 16 · 부록] 데모 영상 목록표. 이 덱에 쓰인 6개 클립(V1~V6)과 각 실행 스크립트/내용/등장 슬라이드.
s = slide()
head(s, 'APPENDIX', 'Demo inventory — every clip in this deck is reproducible')
# (라벨, 스크립트파일, 내용설명, 어느슬라이드/backup, 색) 데모 6종
demos = [
    ('V1', 'demo_wholebuild_anim.jl', 'baseline healthy build, no OOD', 'slide 3', GRAY),
    ('V2', 'demo_respec_replace_anim.jl', 'LLM emits ReplaceAgent → spare hot-swaps in', 'slide 7', ORANGE),
    ('V3', 'demo_rl_replace_anim.jl', 'the trained RL producer on the same fault', 'backup', BLUE),
    ('V4', 'demo_surrogate_stream_anim.jl', 'surrogate scores 5 macros in imagination, OOD stream + live SoC', 'slide 12', TEAL),
    ('V5', 'demo_respec_forbidzone_anim.jl', 'keep-out zone → ForbidZone restage → build completes', 'backup', TEAL),
    ('V6', 'demo_energy_stall_replace_anim.jl', 'SoC → 0 → stall → breakdown → spare replacement', 'backup', BLUE),
]
for i, (k, f, d, where, c) in enumerate(demos):   # 데모 한 줄씩: 색 배지 + 스크립트명 + 설명 + 위치
    y = 2.05 + i * 0.78                             # i번째 행 세로 위치
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.80), Inches(y + 0.03), Inches(0.62), Inches(0.46))  # V1~V6 색 배지
    b.fill.solid(); b.fill.fore_color.rgb = c; b.line.fill.background(); b.shadow.inherit = False
    line(tb(s, 0.80, y + 0.09, 0.62, 0.4), k, 13, WHITE, True, first=True, align=PP_ALIGN.CENTER)   # 배지 안 라벨(V1..)
    line(tb(s, 1.60, y + 0.05, 4.3, 0.45), f, 14, NAVY, True, first=True)     # 스크립트 파일명
    line(tb(s, 6.00, y + 0.07, 5.4, 0.5), d, 13, GRAY, first=True)            # 내용 설명
    line(tb(s, 11.55, y + 0.07, 1.3, 0.4), where, 12, c, True, first=True)    # 등장 위치(슬라이드 번호/backup)
line(tb(s, 0.80, 6.75, 11.8, 0.5),
     'cd ConstructionBots.jl  ·  julia +lts --project=. tools/<script>  →  results/tractor/…/visualization.html',
     13, GRAY, True, first=True)
foot(s)
notes(s, """
Appendix — how to actually produce the videos.

Recipe (same for all six):
  1. cd C:\\Users\\chahj\\PythonCodes\\venv\\ConstructionBots.jl
  2. julia +lts --project=. tools/<script>.jl      (SAVE_ANIM=1 is the default)
  3. open results\\tractor\\greedy_RVO_Dispersion_TangentBug\\visualization.html
  4. screen-record the browser (Win+G / OBS), scrub or let it autoplay, trim to the length noted
  5. save as .mp4 next to the pptx, then in PowerPoint: Insert ▸ Video ▸ This Device, drop it onto the
     dashed orange placeholder, and set Playback ▸ Start: Automatically (+ Loop for V1).

Knobs worth setting before recording V4 (the headline clip):
  OOD_N=4  OOD_KINDS=fault,battery,zone  OOD_SEED=3  SIDEBAR_W=380
  SHRINK / STALL_SOC must stay at 200 / 0.15 — they must match the oracle's DS_SHRINK / DS_STALL,
  or the surrogate on screen is predicting a different world than the one it was trained on.

Suggested cut if the talk over-runs: keep V1 and V4 only. V4 alone carries the thesis.
""")

# 완성된 발표덱을 파일로 저장하고, 만들어진 슬라이드 수를 출력한다.
out = 'ConstructionBots_evolution.pptx'
prs.save(out)   # 현재 폴더에 .pptx 파일 기록
# prs.slides._sldIdLst = 슬라이드 목록의 내부 XML 요소. len(...)으로 총 슬라이드 개수를 셈.
print('wrote', out, '-', len(prs.slides._sldIdLst), 'slides')
